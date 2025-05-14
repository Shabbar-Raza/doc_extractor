from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse
import uvicorn
import os
import tempfile
from typing import Dict, Optional
import io

# Document processing libraries
import PyPDF2
import docx
import pandas as pd
import mammoth
import pptx
from bs4 import BeautifulSoup
import json

app = FastAPI(title="Document Text Extractor")

# Configure CORS to allow requests from anywhere
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class DocumentExtractor:
    """Handles extraction of text from various document formats"""
    
    @staticmethod
    def extract_from_pdf(file_path: str) -> str:
        """Extract text from PDF files"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n\n"
            return text
        except Exception as e:
            return f"Error extracting text from PDF: {str(e)}"

    @staticmethod
    def extract_from_docx(file_path: str) -> str:
        """Extract text from Word documents"""
        try:
            doc = docx.Document(file_path)
            text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            return f"Error extracting text from DOCX: {str(e)}"
    
    @staticmethod
    def extract_from_doc(file_path: str) -> str:
        """Extract text from legacy DOC files using mammoth"""
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.extract_raw_text(docx_file)
                return result.value
        except Exception as e:
            return f"Error extracting text from DOC: {str(e)}"

    @staticmethod
    def extract_from_txt(file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                return file.read()
        except Exception as e:
            return f"Error extracting text from TXT: {str(e)}"

    @staticmethod
    def extract_from_csv(file_path: str) -> str:
        """Extract text from CSV files"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            return f"Error extracting text from CSV: {str(e)}"

    @staticmethod
    def extract_from_excel(file_path: str) -> str:
        """Extract text from Excel files"""
        try:
            df = pd.read_excel(file_path, sheet_name=None)
            result = ""
            for sheet_name, sheet_df in df.items():
                result += f"Sheet: {sheet_name}\n"
                result += sheet_df.to_string() + "\n\n"
            return result
        except Exception as e:
            return f"Error extracting text from Excel: {str(e)}"

    @staticmethod
    def extract_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint files"""
        try:
            pres = pptx.Presentation(file_path)
            text = ""
            
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
                
            return text
        except Exception as e:
            return f"Error extracting text from PowerPoint: {str(e)}"

    @staticmethod
    def extract_from_html(file_path: str) -> str:
        """Extract text from HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                return soup.get_text()
        except Exception as e:
            return f"Error extracting text from HTML: {str(e)}"

    @staticmethod
    def process_file(file_path: str, file_ext: str) -> Dict[str, str]:
        """Process a file based on its extension"""
        try:
            if file_ext == '.pdf':
                text = DocumentExtractor.extract_from_pdf(file_path)
            elif file_ext == '.docx':
                text = DocumentExtractor.extract_from_docx(file_path)
            elif file_ext == '.doc':
                text = DocumentExtractor.extract_from_doc(file_path)
            elif file_ext == '.txt':
                text = DocumentExtractor.extract_from_txt(file_path)
            elif file_ext in ['.csv']:
                text = DocumentExtractor.extract_from_csv(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                text = DocumentExtractor.extract_from_excel(file_path)
            elif file_ext == '.pptx':
                text = DocumentExtractor.extract_from_pptx(file_path)
            elif file_ext in ['.html', '.htm']:
                text = DocumentExtractor.extract_from_html(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']:
                text = "Image processing is not supported in this version."
            else:
                text = f"Unsupported file format: {file_ext}"
            
            return {"text": text}
        except Exception as e:
            return {"error": f"Error processing file: {str(e)}"}


@app.post("/extract-text/")
async def extract_text(file: UploadFile = File(...)) -> Dict[str, str]:
    """
    Extract text from an uploaded document file
    """
    # Check if file is empty
    if not file:
        raise HTTPException(status_code=400, detail="No file provided")
    
    # Create a temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False)
    try:
        # Write the uploaded file to the temporary file
        contents = await file.read()
        with open(temp_file.name, 'wb') as f:
            f.write(contents)
        
        # Get the file extension
        file_ext = os.path.splitext(file.filename)[1].lower()
        
        # Process the file based on its extension
        result = DocumentExtractor.process_file(temp_file.name, file_ext)
        
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
    finally:
        # Clean up the temporary file
        temp_file.close()
        if os.path.exists(temp_file.name):
            os.unlink(temp_file.name)


@app.get("/", response_class=HTMLResponse)
async def read_root():
    """
    Return a simple HTML form for file upload
    """
    html_content = """
    <!DOCTYPE html>
    <html>
    <head>
        <title>Document Text Extractor</title>
        <style>
            body {
                font-family: Arial, sans-serif;
                max-width: 800px;
                margin: 0 auto;
                padding: 20px;
            }
            h1 {
                color: #333;
            }
            form {
                background-color: #f9f9f9;
                padding: 20px;
                border-radius: 8px;
                margin-bottom: 20px;
            }
            input[type="file"] {
                margin: 10px 0;
                padding: 10px;
                width: 100%;
            }
            button {
                background-color: #4285f4;
                color: white;
                border: none;
                padding: 10px 20px;
                border-radius: 4px;
                cursor: pointer;
                font-size: 16px;
            }
            button:hover {
                background-color: #3367d6;
            }
            pre {
                background-color: #f0f0f0;
                padding: 15px;
                border-radius: 4px;
                overflow-x: auto;
                white-space: pre-wrap;
                word-wrap: break-word;
            }
            .result-container {
                margin-top: 20px;
            }
        </style>
    </head>
    <body>
        <h1>Document Text Extractor</h1>
        <form id="upload-form" enctype="multipart/form-data">
            <p>Select a document file to extract text:</p>
            <input type="file" id="file-input" name="file" />
            <button type="submit">Extract Text</button>
        </form>
        
        <div class="result-container">
            <h2>Extracted Text:</h2>
            <pre id="result">Select a file and click "Extract Text" to see the result</pre>
        </div>
        
        <script>
            document.getElementById('upload-form').addEventListener('submit', async (e) => {
                e.preventDefault();
                
                const fileInput = document.getElementById('file-input');
                const resultElement = document.getElementById('result');
                
                if (!fileInput.files.length) {
                    resultElement.textContent = 'Please select a file first';
                    return;
                }
                
                const formData = new FormData();
                formData.append('file', fileInput.files[0]);
                
                resultElement.textContent = 'Processing...';
                
                try {
                    const response = await fetch('/extract-text/', {
                        method: 'POST',
                        body: formData,
                    });
                    
                    if (!response.ok) {
                        throw new Error(`Error: ${response.status} - ${response.statusText}`);
                    }
                    
                    const data = await response.json();
                    
                    if (data.text) {
                        resultElement.textContent = data.text;
                    } else if (data.error) {
                        resultElement.textContent = `Error: ${data.error}`;
                    } else {
                        resultElement.textContent = JSON.stringify(data, null, 2);
                    }
                } catch (error) {
                    resultElement.textContent = `Failed to extract text: ${error.message}`;
                }
            });
        </script>
    </body>
    </html>
    """
    return HTMLResponse(content=html_content)


if __name__ == "__main__":
    # Run the FastAPI app
    uvicorn.run(app, host="0.0.0.0", port=8000)

    from pydantic import BaseModel
import openai
from dotenv import load_dotenv

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

class ChatRequest(BaseModel):
    documentText: str
    userMessage: str

@app.post("/chat/")
async def chat_with_document(request: ChatRequest):
    """
    Accepts extracted document text + user query, returns RAG-style answer using OpenAI.
    """
    system_prompt = (
        "You are an assistant that only answers based on the provided document text.\n"
        "If the information is not in the document, reply: 'The document does not contain information on that topic.'"
    )

    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",  # Or use gpt-4 if available
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Document:\n{request.documentText[:100000]}\n\nUser Question:\n{request.userMessage}"}
            ]
        )
        answer = response['choices'][0]['message']['content']
        return {"answer": answer}
    
    except Exception as e:
        return {"error": str(e)}
